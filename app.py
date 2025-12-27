import random
import sys
import time
import os
import re
import json
import uuid
import base64
import requests
import mimetypes
from datetime import datetime
import tempfile
import speech_recognition as sr
from moviepy import VideoFileClip

import cv2
import numpy as np

from PyQt6.QtWidgets import *
from PyQt6.QtCore import *
from PyQt6.QtGui import *
from PyQt6.QtMultimedia import QMediaPlayer, QAudioOutput
from PyQt6.QtMultimediaWidgets import QVideoWidget

import pypdf
import docx
from pptx import Presentation

DEFAULT_API_KEY = ""
SETTINGS_FILE = "settings.json"
WORKSPACE_FOLDER = "ai_workspace"
MEDIA_FOLDER = os.path.join(WORKSPACE_FOLDER, "images")
CHATS_FILE = "chats.json"

AUTO_CHUNK_THRESHOLD = 30000

MAX_VIDEO_SIZE_BYTES = 20 * 1024 * 1024

#Mdl2 icons cuz windows 11 emojis look ugly
MDL2_ICONS = {
    'attach': '\uE723',      # Paperclip
    'folder': '\uE8B7',      # Folder
    'settings': '\uE713',    # Settings/Gear
    'new': '\uE710',         # Add/Plus
    'search': '\uE721',      # Search
    'file': '\uE8A5',        # Document
    'image': '\uE91B',       # Picture
    'video': '\uE714',       # Video
    'delete': '\uE74D',      # Delete
}



TEXT_MODELS_FAST = [
    "gemini-fast",
    "openai-fast",
    "grok",
    "mistral",
    "qwen-coder",
    "llama",
]

TEXT_MODELS_LARGE = [
    "claude-large",
    "gemini-large",
    "openai-large",
]

TEXT_MODELS_OTHER = [
    "openai",
    "gemini",
    "claude",
    "gemini-search",
    "chickytutor",
    "perplexity-fast",
    "perplexity-reasoning",
    "kimi-k2-thinking",
    "deepseek",
]

IMAGE_MODELS = [
    "flux",
    "zimage",
    "turbo",
    "gptimage",
    "seedream",
    "kontext",
    "nanobanana",
    "seedream-pro",
    "nanobanana-pro",
]

TEXT_MODELS = TEXT_MODELS_FAST + TEXT_MODELS_LARGE + TEXT_MODELS_OTHER

VISION_MODELS = [
    "gemini-fast",
    "openai-fast",
    "gemini-large",
    "openai-large",
    "gemini-search",
    "claude-large",
    "openai",
    "gemini",
    "grok",
    "claude",
    "claude-fast",
    "seedream",
    "kontext",
    "nanobanana",
    "seedream-pro",
    "nanobanana-pro",
    "gptimage",
]


class SettingsManager:
    @staticmethod
    def get_api_key():
        if os.path.exists(SETTINGS_FILE):
            try:
                with open(SETTINGS_FILE, 'r') as f:
                    data = json.load(f)
                return data.get("api_key", DEFAULT_API_KEY)
            except:
                pass
        return DEFAULT_API_KEY

    @staticmethod
    def save_api_key(key):
        data = {"api_key": key}
        with open(SETTINGS_FILE, 'w') as f:
            json.dump(data, f)


SCROLLBAR_STYLESHEET = """
QScrollBar:vertical {
    border: none;
    background: #252525;
    width: 10px;
    margin: 0px 0px 0px 0px;
    border-radius: 0px;
}
QScrollBar::handle:vertical {
    background: #555;
    min-height: 20px;
    border-radius: 5px;
}
QScrollBar::handle:vertical:hover {
    background: #777;
}
QScrollBar::add-line:vertical {
    height: 0px;
    subcontrol-position: bottom;
    subcontrol-origin: margin;
}
QScrollBar::sub-line:vertical {
    height: 0px;
    subcontrol-position: top;
    subcontrol-origin: margin;
}
QScrollBar::add-page:vertical, QScrollBar::sub-page:vertical {
    background: none;
}

QScrollBar:horizontal {
    border: none;
    background: #252525;
    height: 10px;
    margin: 0px 0px 0px 0px;
    border-radius: 0px;
}
QScrollBar::handle:horizontal {
    background: #555;
    min-width: 20px;
    border-radius: 5px;
}
QScrollBar::handle:horizontal:hover {
    background: #777;
}
QScrollBar::add-line:horizontal {
    width: 0px;
}
QScrollBar::sub-line:horizontal {
    width: 0px;
}
QScrollBar::add-page:horizontal, QScrollBar::sub-page:horizontal {
    background: none;
}
"""

DIALOG_STYLESHEET = """
    QDialog {
        background-color: #333333;
        color: white;
    }
    QLabel {
        color: white;
        font-size: 14px;
    }
    QLineEdit, QTextEdit {
        background-color: #444444;
        color: white;
        padding: 5px;
        border: 1px solid #555;
        border-radius: 4px;
    }
    QPushButton {
        background-color: #555555;
        color: white;
        padding: 6px 15px;
        border-radius: 4px;
        border: none;
    }
    QPushButton:hover {
        background-color: #666666;
    }
"""


class FileConverter:
    TEXT_EXTENSIONS = {
        '.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.csv', '.xml',
        '.java', '.c', '.cpp', '.h', '.cs', '.php', '.rb', '.go', '.rs',
        '.swift', '.ts', '.sh', '.bat', '.ps1', '.sql', '.yaml', '.yml',
        '.ini', '.toml', '.cfg', '.log', '.env', '.dockerfile'
    }

    VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv'}
    IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}

    @staticmethod
    def get_file_extension(filepath):
        return os.path.splitext(filepath)[1].lower()

    @staticmethod
    def is_video(filepath):
        return FileConverter.get_file_extension(filepath) in FileConverter.VIDEO_EXTENSIONS

    @staticmethod
    def is_image(filepath):
        return FileConverter.get_file_extension(filepath) in FileConverter.IMAGE_EXTENSIONS

    @staticmethod
    def get_mime_type(filepath):
        mime_type, _ = mimetypes.guess_type(filepath)
        return mime_type or "application/octet-stream"

    @staticmethod
    def read_pdf(filepath):
        try:
            reader = pypdf.PdfReader(filepath)
            text = []
            for page in reader.pages:
                text.append(page.extract_text() or "")
            return "\n".join(text)
        except Exception as e:
            return f"[Error reading PDF: {str(e)}]"

    @staticmethod
    def read_docx(filepath):
        try:
            doc = docx.Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            return f"[Error reading DOCX: {str(e)}]"

    @staticmethod
    def read_pptx(filepath):
        try:
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            return f"[Error reading PPTX: {str(e)}]"

    @staticmethod
    def encode_media_base64(filepath):
        try:
            with open(filepath, "rb") as media_file:
                return base64.b64encode(media_file.read()).decode('utf-8')
        except Exception as e:
            print(f"Error encoding media: {e}")
            return None

    @staticmethod
    def get_video_thumbnail(filepath):
        """Extracts the first frame of a video using OpenCV and returns a QPixmap."""
        try:
            cap = cv2.VideoCapture(filepath)
            if not cap.isOpened():
                return None

            ret, frame = cap.read()
            cap.release()

            if not ret:
                return None

            frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            h, w, ch = frame.shape
            bytes_per_line = ch * w
            q_img = QImage(frame.data, w, h, bytes_per_line, QImage.Format.Format_RGB888)

            return QPixmap.fromImage(q_img)
        except Exception as e:
            print(f"Error extracting video thumbnail: {e}")
            return None

    @staticmethod
    def extract_frames_for_api(filepath, max_frames=20, quality=70):
        """
        Extracts keyframes from a video to simulate 'watching' it.
        """
        try:
            cap = cv2.VideoCapture(filepath)
            if not cap.isOpened():
                return []

            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            step = max(1, total_frames // max_frames)

            frames_b64 = []
            count = 0

            while cap.isOpened():
                ret, frame = cap.read()
                if not ret:
                    break

                if count % step == 0:
                    h, w = frame.shape[:2]
                    scale = min(512 / w, 512 / h)
                    new_w, new_h = int(w * scale), int(h * scale)
                    resized = cv2.resize(frame, (new_w, new_h))

                    _, buffer = cv2.imencode('.jpg', resized, [int(cv2.IMWRITE_JPEG_QUALITY), quality])
                    b64_str = base64.b64encode(buffer).decode('utf-8')
                    frames_b64.append(b64_str)

                    if len(frames_b64) >= max_frames:
                        break

                count += 1

            cap.release()
            return frames_b64
        except Exception as e:
            print(f"Error extracting frames: {e}")
            return []

    @staticmethod
    def extract_audio_text_from_video(filepath):
        video_clip = None
        temp_audio_path = None
        try:
            print(f"[Info] Extracting audio from {filepath}...")
            video_clip = VideoFileClip(filepath)

            if not video_clip.audio:
                return "[Video has no audio track]"

            duration = getattr(video_clip, 'duration', None)

            if duration and duration > 0:
                max_dur = min(duration, 90.0)
                try:
                    audio_subclip = video_clip.subclip(0, max_dur).audio
                except Exception:
                    audio_subclip = video_clip.audio.with_end(max_dur)
            else:
                audio_subclip = video_clip.audio

            fd, temp_audio_path = tempfile.mkstemp(suffix=".wav")
            os.close(fd)

            audio_subclip.write_audiofile(
                temp_audio_path,
                codec="pcm_s16le",
                fps=16000,
                ffmpeg_params=["-ac", "1"]
            )

            recognizer = sr.Recognizer()
            with sr.AudioFile(temp_audio_path) as source:
                audio_data = recognizer.record(source)
                return recognizer.recognize_google(audio_data)

        except sr.UnknownValueError:
            return "[Audio was present but unintelligible]"
        except Exception as e:
            print(f"[Error] Audio processing failed: {e}")
            return f"[Error extracting audio: {str(e)}]"
        finally:
            if video_clip:
                video_clip.close()
            if temp_audio_path and os.path.exists(temp_audio_path):
                try:
                    os.remove(temp_audio_path)
                except:
                    pass

    @staticmethod
    def convert(filepath):
        ext = FileConverter.get_file_extension(filepath)

        if ext == ".pdf":
            return FileConverter.read_pdf(filepath)
        elif ext in [".docx", ".doc"]:
            return FileConverter.read_docx(filepath)
        elif ext in [".pptx", ".ppt"]:
            return FileConverter.read_pptx(filepath)
        elif ext in FileConverter.TEXT_EXTENSIONS or ext == "":
            try:
                with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
            except Exception as e:
                return f"[Error reading file: {str(e)}]"
        else:
            try:
                with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()
            except:
                return "[Error: Binary or unsupported file format]"


DEFAULT_SYSTEM_PROMPT = f"""
You are an advanced AI assistant with file system access. 
Sometimes files are sent to u in a split variant but that's okay. User sees it normally so u don't need to reconstruct anything.
You don't have a main purpose of only working with files (U should also be a good chatbot).
You are working inside an isolated project folder named '{WORKSPACE_FOLDER}'.
You can see attached images and videos.

PROJECT MANAGMENT CAPABILITIES:
1. You can CREATE or OVERWRITE files.
2. You can DELETE files.
3. You can READ files (if the user attaches them).
4. You can SEE images and WATCH videos if attached.

COMMAND SYNTAX:
To perform actions, you MUST use the following specific block formats in your response.

1. TO CREATE or OVERWRITE A FILE:
:::create file="filename.ext":::
[File Content Here]
:::end_create::: // DO NOT FORGET TO USE :::end_create:::

2. TO DELETE A FILE:
:::delete file="filename.ext":::

RULES:
- Do not use markdown code blocks inside the :::create::: tags. Just write the raw text.
- Only modify files if explicitly asked or if it solves the user's coding problem.
- If you create a file, tell the user you have done so in your text response.
- Don't spend a lot of time in thinking mode.
"""


class FileManager:
    @staticmethod
    def ensure_workspace():
        if not os.path.exists(WORKSPACE_FOLDER):
            os.makedirs(WORKSPACE_FOLDER)
        if not os.path.exists(MEDIA_FOLDER):
            os.makedirs(MEDIA_FOLDER)

    @staticmethod
    def is_safe_path(filename):
        base_dir = os.path.abspath(WORKSPACE_FOLDER)
        target_path = os.path.abspath(os.path.join(WORKSPACE_FOLDER, filename))
        return target_path.startswith(base_dir)

    @staticmethod
    def write_file(filename, content):
        if not FileManager.is_safe_path(filename):
            return "Error: unsafe path."

        path = os.path.join(WORKSPACE_FOLDER, filename)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(content)
        return f"File '{filename}' created/updated successfully."

    @staticmethod
    def delete_file(filename):
        if not FileManager.is_safe_path(filename):
            return "Error: unsafe path."

        path = os.path.join(WORKSPACE_FOLDER, filename)
        if os.path.exists(path):
            os.remove(path)
            return f"File '{filename}' deleted."
        return f"File '{filename}' not found."

    @staticmethod
    def save_media_from_bytes(media_bytes, extension=".png"):
        try:
            image_id = str(uuid.uuid4())
            filename = f"{image_id}{extension}"
            path = os.path.join(MEDIA_FOLDER, filename)

            with open(path, "wb") as f:
                f.write(media_bytes)

            return f"images/{filename}"
        except Exception as e:
            print(f"Error saving media: {e}")
            return None


class ChatStorage:
    @staticmethod
    def load_chats():
        if not os.path.exists(CHATS_FILE):
            return {}
        try:
            with open(CHATS_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return {}

    @staticmethod
    def save_chats(chats):
        with open(CHATS_FILE, 'w', encoding='utf-8') as f:
            json.dump(chats, f, indent=4)


class AIResponseThread(QThread):
    response_chunk = pyqtSignal(str)
    response_complete = pyqtSignal(str)
    error_occurred = pyqtSignal(str)
    image_generated = pyqtSignal(bytes, str)

    def __init__(self, model, messages, is_image_model=False):
        super().__init__()
        self.model = model
        self.messages = messages
        self.is_image_model = is_image_model
        self._is_running = True
        self.api_key = SettingsManager.get_api_key()

    def stop(self):
        self._is_running = False

    def run(self):
        if self.is_image_model:
            self.generate_image()
        else:
            self.generate_text()

    def generate_image(self):
        try:
            prompt = ""
            for msg in reversed(self.messages):
                if msg["role"] == "user":
                    content = msg["content"]
                    if isinstance(content, str):
                        prompt = content
                    elif isinstance(content, list):
                        for part in content:
                            if part.get("type") == "text":
                                prompt = part.get("text", "")
                                break
                    break

            if not prompt:
                self.error_occurred.emit("No prompt found for image generation")
                return

            url = f"https://image.pollinations.ai/prompt/{requests.utils.quote(prompt)}"

            headers = {
                "Authorization": f"Bearer {self.api_key}"
            }

            params = {
                "model": self.model,
                "seed": random.randint(1, 10000),
                "width": 1024,
                "height": 1024,
                "private": "true",
                "safe": "false",
                "nologo": "true"
            }

            response = requests.get(url, headers=headers, params=params, timeout=250)
            response.raise_for_status()

            if response.content:
                saved_path = FileManager.save_media_from_bytes(response.content, ".png")
                if saved_path:
                    self.image_generated.emit(response.content, saved_path)
                    self.response_complete.emit(f"🖼️ Generated image using {self.model}")
                else:
                    self.error_occurred.emit("Failed to save generated image")
            else:
                self.error_occurred.emit("No image data received")

        except Exception as e:
            self.error_occurred.emit(f"Image generation failed: {str(e)}")

    def generate_text(self):
        max_retries = 6
        base_delay = 1.0

        for attempt in range(max_retries + 1):
            if not self._is_running: return
            try:
                full_response = ""
                prepared_messages = self.prepare_messages_for_api(self.messages)

                url = "https://gen.pollinations.ai/v1/chat/completions"

                headers = {
                    "Authorization": f"Bearer {self.api_key}",
                    "Content-Type": "application/json"
                }

                payload = {
                    "model": self.model,
                    "messages": prepared_messages,
                    "stream": True,
                    "temperature": 1.0
                }

                response = requests.post(url, headers=headers, json=payload, stream=True, timeout=250)
                response.raise_for_status()

                for line in response.iter_lines():
                    if not self._is_running: return
                    if line:
                        line_str = line.decode('utf-8').strip()
                        if line_str.startswith('data: '):
                            json_str = line_str[6:]
                            if json_str == '[DONE]':
                                break
                            try:
                                chunk_json = json.loads(json_str)

                                choices = chunk_json.get("choices", [])
                                if choices:
                                    content = choices[0].get("delta", {}).get("content", "")
                                    if content:
                                        self.response_chunk.emit(content)
                                        full_response += content
                                        self.msleep(10)
                            except Exception:
                                continue

                if self._is_running:
                    self.response_complete.emit(full_response.strip())
                return

            except Exception as e:
                error_str = str(e).lower()
                retryable = (
                        "443" in error_str or
                        "500" in error_str or
                        "502" in error_str or
                        "503" in error_str or
                        "504" in error_str or
                        "bad gateway" in error_str or
                        "service unavailable" in error_str or
                        "timeout" in error_str or
                        "connection" in error_str or
                        "reset" in error_str or
                        "rate limit" in error_str or
                        "413" in error_str or
                        "429" in error_str
                )

                if attempt < max_retries and self._is_running and retryable:
                    delay = base_delay * (2 ** attempt) + (0.1 * attempt)
                    print(f"\n[Retry] {e} → Retrying in {delay:.1f}s...")
                    time.sleep(delay)
                    continue
                else:
                    final_error = f"Failed after {attempt} retries: {str(e)}"
                    print("\n[Error]", final_error)
                    self.error_occurred.emit(final_error)
                    return

    def prepare_messages_for_api(self, messages):
        """
        Scan messages for local file paths.
        - Images: Convert to base64.
        - Videos:
            1. Extract Audio Transcript so AI 'hears' it.
            2. If it's the LATEST message: Send visuals (Base64 or Frames).
            3. If it's an OLD message: Remove visuals (save tokens), keep Transcript.
        """
        api_messages = []
        total_msgs = len(messages)

        for idx, msg in enumerate(messages):
            is_latest = (idx == total_msgs - 1)
            new_msg = msg.copy()

            if isinstance(new_msg.get("content"), list):
                new_content = []
                for part in new_msg["content"]:
                    if part.get("type") == "image_url":
                        url = part["image_url"]["url"]
                        if not url.startswith("data:") and not url.startswith("http"):
                            full_path = os.path.join(WORKSPACE_FOLDER, url)
                            if os.path.exists(full_path):
                                b64 = FileConverter.encode_media_base64(full_path)
                                if b64:
                                    part_copy = part.copy()
                                    part_copy["image_url"] = {"url": f"data:image/png;base64,{b64}"}
                                    new_content.append(part_copy)
                                    continue
                        new_content.append(part)

                    elif part.get("type") == "video_url":
                        url = part["video_url"]["url"]
                        if not url.startswith("data:") and not url.startswith("http"):
                            full_path = os.path.join(WORKSPACE_FOLDER, url)
                            if os.path.exists(full_path):

                                audio_text = FileConverter.extract_audio_text_from_video(full_path)
                                new_content.append({
                                    "type": "text",
                                    "text": f"\n[Transcript of words spoken in video '{os.path.basename(url)}': \"{audio_text}\"]\n"
                                })

                                if is_latest:
                                    file_size = os.path.getsize(full_path)

                                    if file_size < MAX_VIDEO_SIZE_BYTES:
                                        b64 = FileConverter.encode_media_base64(full_path)
                                        mime = FileConverter.get_mime_type(full_path)
                                        if b64:
                                            part_copy = part.copy()
                                            part_copy["video_url"] = {"url": f"data:{mime};base64,{b64}"}
                                            new_content.append(part_copy)
                                            continue

                                    else:
                                        print(
                                            f"[Info] Video too large ({file_size / 1024 / 1024:.2f}MB). Extracting frames...")
                                        frames = FileConverter.extract_frames_for_api(full_path, max_frames=20)
                                        if frames:
                                            new_content.append({"type": "text",
                                                                "text": "[Video content represented by the following keyframes:]"})
                                            for frame_b64 in frames:
                                                new_content.append({
                                                    "type": "image_url",
                                                    "image_url": {
                                                        "url": f"data:image/jpeg;base64,{frame_b64}"
                                                    }
                                                })
                                            continue
                                else:
                                    new_content.append({
                                        "type": "text",
                                        "text": f"[Visual video data for '{os.path.basename(url)}' removed to save resources. See transcript above.]"
                                    })
                                    continue

                        new_content.append(part)
                    else:
                        new_content.append(part)
                new_msg["content"] = new_content
            api_messages.append(new_msg)
        return api_messages


class ClickableLabel(QLabel):
    clicked = pyqtSignal()

    def mousePressEvent(self, event):
        self.clicked.emit()
        super().mousePressEvent(event)


class ClickableImageLabel(QLabel):
    clicked = pyqtSignal()

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            self.clicked.emit()
        super().mousePressEvent(event)


class StreamingMessageBubble(QFrame):
    def __init__(self, sender="assistant", text=""):
        super().__init__()
        self.sender = sender
        self.full_text = text
        self.display_buffer = text

        self.setFrameShape(QFrame.Shape.NoFrame)
        self.setStyleSheet(self._bubble_style())
        self.setSizePolicy(QSizePolicy.Policy.Preferred, QSizePolicy.Policy.Minimum)

        self.layout = QVBoxLayout(self)
        self.layout.setContentsMargins(12, 12, 12, 12)
        self.layout.setSpacing(5)

        self.label = ClickableLabel()
        self.label.setTextFormat(Qt.TextFormat.MarkdownText)
        self.label.setWordWrap(True)
        self.label.setTextInteractionFlags(
            Qt.TextInteractionFlag.TextSelectableByMouse | Qt.TextInteractionFlag.LinksAccessibleByMouse)
        self.label.setFont(QFont("Segoe UI", 11))
        self.label.setStyleSheet("color: white; border: none;")
        self.label.setAlignment(Qt.AlignmentFlag.AlignTop | Qt.AlignmentFlag.AlignLeft)
        self.label.setSizePolicy(QSizePolicy.Policy.Preferred, QSizePolicy.Policy.Minimum)
        self.label.setText(text)
        self.label.linkActivated.connect(self.handle_link)

        self.layout.addWidget(self.label)

        self.typing_indicator = QLabel("...")
        self.typing_indicator.setStyleSheet("color: rgba(255, 255, 255, 0.5); font-style: italic;")
        self.typing_indicator.setVisible(False)
        self.layout.addWidget(self.typing_indicator)

        self.setMaximumWidth(650)

        self.update_timer = QTimer()
        self.update_timer.setInterval(30)
        self.update_timer.timeout.connect(self.update_display_text)
        self.text_dirty = False

    def handle_link(self, url):
        pass

    def resizeEvent(self, event):
        super().resizeEvent(event)

    def _bubble_style(self):
        if self.sender == "user":
            color = "#32CD32"
        elif self.sender == "system":
            color = "#555555"
        else:
            color = "#444444"
        return f"QFrame {{ background-color: {color}; border-radius: 12px; }}"

    def start_streaming(self):
        self.typing_indicator.setVisible(True)
        self.update_timer.start()

    def stop_streaming(self):
        self.typing_indicator.setVisible(False)
        self.update_timer.stop()
        self.update_display_text()

    def add_text_chunk(self, chunk):
        self.full_text += chunk
        self.text_dirty = True

    def set_complete_text(self, text):
        self.full_text = text
        self.stop_streaming()
        self.update_display_text()

    def update_display_text(self):
        if self.text_dirty or self.label.text() != self.full_text:
            self.label.setText(self.full_text)
            self.text_dirty = False


class FileViewerDialog(QDialog):
    def __init__(self, filename, content=None, parent=None):
        super().__init__(parent)
        self.setWindowTitle(f"Viewing: {filename}")
        self.resize(800, 600)
        self.setStyleSheet(DIALOG_STYLESHEET)

        layout = QVBoxLayout(self)

        self.text_edit = QTextEdit()
        self.text_edit.setReadOnly(True)
        self.text_edit.setFont(QFont("Consolas", 10))

        if content is not None:
            self.text_edit.setPlainText(content)
        else:
            path = os.path.join(WORKSPACE_FOLDER, filename)
            if os.path.exists(path):
                try:
                    with open(path, 'r', encoding='utf-8') as f:
                        file_content = f.read()
                    self.text_edit.setPlainText(file_content)
                except Exception as e:
                    self.text_edit.setPlainText(f"Error reading file: {e}")
            else:
                self.text_edit.setPlainText(
                    f"File '{filename}' not found on disk and content not found in chat history.")

        layout.addWidget(self.text_edit)

        close_btn = QPushButton("Close")
        close_btn.clicked.connect(self.accept)
        layout.addWidget(close_btn)


class ImageViewerDialog(QDialog):
    def __init__(self, pixmap, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Image Viewer")
        self.setStyleSheet("background-color: #222;")

        self.original_pixmap = pixmap
        self.scale_factor = 1.0

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)

        self.image_label = QLabel()
        self.image_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.image_label.setPixmap(self.original_pixmap)
        layout.addWidget(self.image_label)

        screen = QApplication.primaryScreen().availableGeometry()

        img_w = pixmap.width()
        img_h = pixmap.height()

        max_w = screen.width() - 100
        max_h = screen.height() - 100

        final_w = min(img_w, max_w)
        final_h = min(img_h, max_h)

        self.setFixedSize(final_w, final_h)

        self.setWindowFlags(self.windowFlags() & ~Qt.WindowType.WindowMaximizeButtonHint)

    def wheelEvent(self, event):
        if event.angleDelta().y() > 0:
            self.zoom(1.1)
        else:
            self.zoom(0.9)

    def zoom(self, factor):
        self.scale_factor *= factor
        self.scale_factor = max(0.1, min(self.scale_factor, 5.0))

        new_w = int(self.original_pixmap.width() * self.scale_factor)
        new_h = int(self.original_pixmap.height() * self.scale_factor)

        scaled = self.original_pixmap.scaled(
            new_w,
            new_h,
            Qt.AspectRatioMode.KeepAspectRatio,
            Qt.TransformationMode.SmoothTransformation
        )

        self.image_label.setPixmap(scaled)


class ClickableVideoWidget(QVideoWidget):
    """Custom VideoWidget to handle mouse clicks for Play/Pause"""
    clicked = pyqtSignal()

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            self.clicked.emit()
        super().mousePressEvent(event)


class VideoPlayerDialog(QDialog):
    def __init__(self, video_path, parent=None):
        super().__init__(parent)
        self.setWindowTitle(f"Media Player - {os.path.basename(video_path)}")
        self.resize(900, 650)
        self.setStyleSheet("""
            QDialog { background-color: #1e1e1e; color: white; }
            QLabel { color: #ddd; font-family: 'Segoe UI', sans-serif; }
            QToolTip { background-color: #333; color: white; border: 1px solid #555; }
        """)

        self.layout = QVBoxLayout(self)
        self.layout.setContentsMargins(0, 0, 0, 0)
        self.layout.setSpacing(0)

        self.video_container = QWidget()
        self.video_container.setStyleSheet("background-color: black;")
        video_layout = QVBoxLayout(self.video_container)
        video_layout.setContentsMargins(0, 0, 0, 0)

        self.video_widget = ClickableVideoWidget()
        self.video_widget.clicked.connect(self.toggle_playback)
        video_layout.addWidget(self.video_widget)
        self.layout.addWidget(self.video_container, stretch=1)

        self.controls_widget = QWidget()
        self.controls_widget.setStyleSheet("""
            QWidget { background-color: #2b2b2b; border-top: 1px solid #3d3d3d; }
            QPushButton { 
                background-color: transparent; 
                border: none; 
                border-radius: 4px; 
                padding: 4px;
                color: #ddd;
            }
            QPushButton:hover { background-color: #3d3d3d; }
        """)
        self.controls_layout = QVBoxLayout(self.controls_widget)
        self.controls_layout.setContentsMargins(15, 10, 15, 15)

        self.slider = QSlider(Qt.Orientation.Horizontal)
        self.slider.setCursor(Qt.CursorShape.PointingHandCursor)
        self.slider.setStyleSheet("""
            QSlider::groove:horizontal { 
                border: 1px solid #3d3d3d; 
                height: 6px; 
                background: #1a1a1a; 
                border-radius: 3px; 
            }
            QSlider::handle:horizontal { 
                background: #32CD32; 
                width: 16px; 
                height: 16px; 
                margin: -6px 0; 
                border-radius: 8px; 
            }
            QSlider::handle:horizontal:hover { background: #40E0D0; }
            QSlider::sub-page:horizontal { background: #32CD32; border-radius: 3px; }
        """)
        self.slider.sliderMoved.connect(self.set_position)
        self.slider.sliderPressed.connect(self.slider_pressed)
        self.slider.sliderReleased.connect(self.slider_released)
        self.controls_layout.addWidget(self.slider)

        h_controls = QHBoxLayout()
        h_controls.setSpacing(15)

        self.play_btn = QPushButton("\uE768")
        self.play_btn.setFixedSize(32, 32)
        self.play_btn.setFont(QFont("Segoe MDL2 Assets", 16))
        self.play_btn.setStyleSheet("font-weight: normal;")
        self.play_btn.clicked.connect(self.toggle_playback)
        h_controls.addWidget(self.play_btn)

        self.time_label = QLabel("00:00 / 00:00")
        self.time_label.setStyleSheet("font-family: monospace; font-size: 12px; color: #aaa;")
        h_controls.addWidget(self.time_label)

        h_controls.addStretch()

        self.speed_combo = QComboBox()
        self.speed_combo.addItems(["0.5x", "1.0x", "1.5x", "2.0x"])
        self.speed_combo.setCurrentIndex(1)
        self.speed_combo.setFixedWidth(80)
        self.speed_combo.setCursor(Qt.CursorShape.PointingHandCursor)
        self.speed_combo.setStyleSheet("""
            QComboBox {
                background-color: #404040;
                color: #ffffff;
                border: 1px solid #555;
                border-radius: 6px;
                padding: 6px 15px;
                font-size: 13px;
            }
            QComboBox:hover {
                background-color: #4a4a4a;
                border: 1px solid #32CD32;
            }
            QComboBox::drop-down {
                border: none;
                width: 20px;
            }
            QComboBox::down-arrow {
                image: none;
                border-left: 5px solid transparent;
                border-right: 5px solid transparent;
                border-top: 6px solid #aaaaaa;
                margin-right: 8px;
            }
            QComboBox QAbstractItemView {
                background-color: #353535;
                color: white;
                border: 1px solid #32CD32;
                outline: none;
            }
            QComboBox QAbstractItemView::item {
                padding: 5px;
                min-height: 25px;
                outline: none;
            }
            QComboBox QAbstractItemView::item:selected,
            QComboBox QAbstractItemView::item:focus {
                background-color: #32CD32;
                color: white;
            }
        """)
        self.speed_combo.currentTextChanged.connect(self.change_speed)
        h_controls.addWidget(self.speed_combo)

        self.vol_icon = QPushButton("\uE767")
        self.vol_icon.setFont(QFont("Segoe MDL2 Assets", 16))
        self.vol_icon.clicked.connect(self.toggle_mute)
        self.vol_icon.setFixedSize(30, 30)
        h_controls.addWidget(self.vol_icon)

        self.volume_slider = QSlider(Qt.Orientation.Horizontal)
        self.volume_slider.setFixedWidth(100)
        self.volume_slider.setRange(0, 100)
        self.volume_slider.setValue(70)
        self.volume_slider.setStyleSheet("""
            QSlider::groove:horizontal { height: 4px; background: #555; border-radius: 2px; }
            QSlider::handle:horizontal { background: #ddd; width: 12px; height: 12px; margin: -4px 0; border-radius: 6px; }
            QSlider::sub-page:horizontal { background: #ddd; border-radius: 2px; }
        """)
        self.volume_slider.valueChanged.connect(self.set_volume)
        h_controls.addWidget(self.volume_slider)

        self.controls_layout.addLayout(h_controls)
        self.layout.addWidget(self.controls_widget)

        self.player = QMediaPlayer()
        self.audio_output = QAudioOutput()
        self.player.setAudioOutput(self.audio_output)
        self.player.setVideoOutput(self.video_widget)

        self.player.positionChanged.connect(self.position_changed)
        self.player.durationChanged.connect(self.duration_changed)
        self.player.mediaStatusChanged.connect(self.media_status_changed)
        self.player.errorOccurred.connect(self.handle_errors)

        self.audio_output.setVolume(0.7)
        self.player.setSource(QUrl.fromLocalFile(video_path))
        self.player.play()
        self.play_btn.setText("\uE769")

        self.is_slider_dragged = False
        self.duration = 0

    def keyPressEvent(self, event):
        """Keyboard Shortcuts"""
        if event.key() == Qt.Key.Key_Space:
            self.toggle_playback()
        elif event.key() == Qt.Key.Key_Right:
            self.player.setPosition(self.player.position() + 5000)
        elif event.key() == Qt.Key.Key_Left:
            self.player.setPosition(self.player.position() - 5000)
        elif event.key() == Qt.Key.Key_Up:
            self.volume_slider.setValue(self.volume_slider.value() + 5)
        elif event.key() == Qt.Key.Key_Down:
            self.volume_slider.setValue(self.volume_slider.value() - 5)
        elif event.key() == Qt.Key.Key_Escape:
            self.close()
        else:
            super().keyPressEvent(event)

    def toggle_playback(self):
        if self.player.playbackState() == QMediaPlayer.PlaybackState.PlayingState:
            self.player.pause()
            self.play_btn.setText("\uE768")
        else:
            self.player.play()
            self.play_btn.setText("\uE769")

    def toggle_mute(self):
        if self.audio_output.isMuted():
            self.audio_output.setMuted(False)
            self.vol_icon.setText("\uE767")
            self.volume_slider.setEnabled(True)
        else:
            self.audio_output.setMuted(True)
            self.vol_icon.setText("\uE74F")
            self.volume_slider.setEnabled(False)

    def change_speed(self, text):
        rate = float(text.replace("x", ""))
        self.player.setPlaybackRate(rate)

    def set_volume(self, value):
        self.audio_output.setVolume(value / 100)
        if value == 0:
            self.vol_icon.setText("\uE74F")
        else:
            self.vol_icon.setText("\uE767")

    def media_status_changed(self, status):
        if status == QMediaPlayer.MediaStatus.EndOfMedia:
            self.play_btn.setText("\uE768")

    def slider_pressed(self):
        self.is_slider_dragged = True

    def slider_released(self):
        self.is_slider_dragged = False
        self.player.setPosition(self.slider.value())

    def set_position(self, position):
        if self.is_slider_dragged:
            self.update_time_label(position)

    def position_changed(self, position):
        if not self.is_slider_dragged:
            self.slider.setValue(position)
            self.update_time_label(position)

    def duration_changed(self, duration):
        self.duration = duration
        self.slider.setRange(0, duration)
        self.update_time_label(self.player.position())

    def update_time_label(self, position):
        def fmt(ms):
            seconds = (ms // 1000) % 60
            minutes = (ms // 60000)
            return f"{minutes:02}:{seconds:02}"

        self.time_label.setText(f"{fmt(position)} / {fmt(self.duration)}")

    def handle_errors(self):
        self.play_btn.setEnabled(False)
        self.time_label.setText("Error: Could not load media")

    def closeEvent(self, event):
        self.player.stop()
        self.player.setSource(QUrl())
        self.player.setVideoOutput(None)
        self.player.deleteLater()
        self.audio_output.deleteLater()
        super().closeEvent(event)


class PromptEditorDialog(QDialog):
    def __init__(self, current_prompt, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Edit System Prompt")
        self.resize(600, 400)
        self.setStyleSheet(DIALOG_STYLESHEET)

        layout = QVBoxLayout(self)

        label = QLabel("System Prompt for this Chat:")
        layout.addWidget(label)

        self.text_edit = QTextEdit()
        self.text_edit.setPlainText(current_prompt)
        layout.addWidget(self.text_edit)

        btn_layout = QHBoxLayout()
        save_btn = QPushButton("Save")
        save_btn.clicked.connect(self.accept)
        cancel_btn = QPushButton("Cancel")
        cancel_btn.clicked.connect(self.reject)

        btn_layout.addStretch()
        btn_layout.addWidget(cancel_btn)
        btn_layout.addWidget(save_btn)
        layout.addLayout(btn_layout)

    def get_prompt(self):
        return self.text_edit.toPlainText()


class SettingsDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Settings")
        self.setStyleSheet(DIALOG_STYLESHEET)
        layout = QVBoxLayout(self)
        layout.addWidget(QLabel("Pollinations.AI API Key:"))
        self.key_input = QLineEdit()
        self.key_input.setText(SettingsManager.get_api_key())
        layout.addWidget(self.key_input)
        btn_row = QHBoxLayout()
        reset_btn = QPushButton("Reset to Default")
        reset_btn.clicked.connect(lambda: self.key_input.setText(DEFAULT_API_KEY))
        save_btn = QPushButton("Save")
        save_btn.clicked.connect(self.save_and_close)
        btn_row.addWidget(reset_btn)
        btn_row.addWidget(save_btn)
        layout.addLayout(btn_row)

    def save_and_close(self):
        SettingsManager.save_api_key(self.key_input.text())
        self.accept()


class AIChatApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("AI Chat Studio")
        self.setGeometry(100, 100, 1100, 750)
        FileManager.ensure_workspace()

        self.setStyleSheet(SCROLLBAR_STYLESHEET)

        palette = QPalette()
        palette.setColor(QPalette.ColorRole.Window, QColor(30, 30, 30))
        palette.setColor(QPalette.ColorRole.WindowText, Qt.GlobalColor.white)
        self.setPalette(palette)

        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_h_layout = QHBoxLayout(central_widget)
        main_h_layout.setContentsMargins(0, 0, 0, 0)
        main_h_layout.setSpacing(0)

        self.splitter = QSplitter(Qt.Orientation.Horizontal)
        self.splitter.setHandleWidth(2)
        self.splitter.setStyleSheet("QSplitter::handle { background-color: #444; }")

        self.sidebar_widget = QWidget()
        self.sidebar_widget.setStyleSheet("background-color: #202020;")
        self.sidebar_widget.setMinimumWidth(200)
        self.sidebar_widget.setMaximumWidth(350)
        sidebar_layout = QVBoxLayout(self.sidebar_widget)
        sidebar_layout.setContentsMargins(10, 10, 10, 10)

        self.new_chat_btn = QPushButton("+ New Chat")
        self.new_chat_btn.setStyleSheet("""
            QPushButton { background-color: #32CD32; color: white; padding: 10px; border-radius: 5px; font-weight: bold; }
            QPushButton:hover { background-color: #2ECC71; }
        """)
        self.new_chat_btn.clicked.connect(self.create_new_chat)
        sidebar_layout.addWidget(self.new_chat_btn)

        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText("🔎︎ Search chats...")
        self.search_input.setStyleSheet("""
            QLineEdit { background-color: #333; color: white; border: 1px solid #444; border-radius: 5px; padding: 5px; }
            QLineEdit:focus { border: 1px solid #32CD32; }
        """)
        self.search_input.textChanged.connect(self.filter_chats)
        sidebar_layout.addWidget(self.search_input)

        self.chat_list_widget = QListWidget()
        self.chat_list_widget.setStyleSheet("""
            QListWidget { background-color: #202020; border: none; outline: none; }
            QListWidget::item { color: #ccc; padding: 10px; border-radius: 5px; }
            QListWidget::item:hover { background-color: #333; }
            QListWidget::item:selected { background-color: #444; color: white; border-left: 3px solid #32CD32; }
        """)
        self.chat_list_widget.itemClicked.connect(self.load_selected_chat)

        self.chat_list_widget.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        self.chat_list_widget.customContextMenuRequested.connect(self.show_context_menu)

        sidebar_layout.addWidget(self.chat_list_widget)

        self.right_widget = QWidget()
        right_layout = QVBoxLayout(self.right_widget)

        self.setAcceptDrops(True)

        header_layout = QHBoxLayout()

        self.model_type_combo = QComboBox()
        self.model_type_combo.addItems(["Text Models", "Image Models"])
        self.model_type_combo.currentTextChanged.connect(self.on_model_type_changed)
        self.model_type_combo.setCursor(Qt.CursorShape.PointingHandCursor)
        self.model_type_combo.setStyleSheet("""
            QComboBox {
                background-color: #404040;
                color: #ffffff;
                border: 1px solid #555;
                border-radius: 6px;
                padding: 6px 15px;
                min-width: 120px;
                font-size: 13px;
            }
            QComboBox:hover {
                background-color: #4a4a4a;
                border: 1px solid #32CD32;
            }
            QComboBox::drop-down {
                border: none;
                width: 20px;
            }
            QComboBox::down-arrow {
                image: none;
                border-left: 5px solid transparent;
                border-right: 5px solid transparent;
                border-top: 6px solid #aaaaaa;
                margin-right: 8px;
            }
            QComboBox QAbstractItemView {
                background-color: #353535;
                color: white;
                border: 1px solid #32CD32;
                outline: none;
            }
            QComboBox QAbstractItemView::item {
                padding: 5px;
                min-height: 25px;
                outline: none;
            }
            QComboBox QAbstractItemView::item:selected,
            QComboBox QAbstractItemView::item:focus {
                background-color: #32CD32;
                color: white;
            }
        """)

        self.model_combo = QComboBox()
        self.model_combo.addItems(TEXT_MODELS)
        self.model_combo.setCursor(Qt.CursorShape.PointingHandCursor)
        self.model_combo.setStyleSheet("""
                    QComboBox {
                        background-color: #404040;
                        color: #ffffff;
                        border: 1px solid #555;
                        border-radius: 6px;
                        padding: 6px 15px;
                        min-width: 120px;
                        font-size: 13px;
                    }
                    QComboBox:hover {
                        background-color: #4a4a4a;
                        border: 1px solid #32CD32;
                    }
                    QComboBox::drop-down {
                        border: none;
                        width: 20px;
                    }
                    QComboBox::down-arrow {
                        image: none;
                        border-left: 5px solid transparent;
                        border-right: 5px solid transparent;
                        border-top: 6px solid #aaaaaa;
                        margin-right: 8px;
                    }
                    QComboBox QAbstractItemView {
                        background-color: #353535;
                        color: white;
                        border: 1px solid #32CD32;
                        outline: none;
                    }
                    QComboBox QAbstractItemView::item {
                        padding: 5px;
                        min-height: 25px;
                        outline: none;
                    }
                    QComboBox QAbstractItemView::item:selected,
                    QComboBox QAbstractItemView::item:focus {
                        background-color: #32CD32;
                        color: white;
                    }
                """)

        self.prompt_btn = QPushButton("⚙️ Prompt")
        self.prompt_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self.prompt_btn.clicked.connect(self.open_prompt_editor)
        self.prompt_btn.setStyleSheet(self.header_btn_style())

        self.workspace_btn = QPushButton(f"🗁 {WORKSPACE_FOLDER}")
        self.workspace_btn.clicked.connect(self.open_workspace_folder)
        self.workspace_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self.workspace_btn.setStyleSheet(self.header_btn_style())

        self.settings_btn = QPushButton("⚙️ Settings")
        self.settings_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self.settings_btn.clicked.connect(self.open_settings)
        self.settings_btn.setStyleSheet(self.header_btn_style())

        header_layout.addWidget(self.model_type_combo)
        header_layout.addWidget(self.model_combo)
        header_layout.addWidget(self.prompt_btn)
        header_layout.addStretch()
        header_layout.addWidget(self.workspace_btn)
        header_layout.addWidget(self.settings_btn)
        right_layout.addLayout(header_layout)

        self.chat_area = QScrollArea()
        self.chat_area.setWidgetResizable(True)
        self.chat_area.setStyleSheet("border: none; background-color: #252525;")

        self.chat_widget_container = QWidget()
        self.chat_widget_container.setStyleSheet("background-color: #252525;")
        self.chat_layout = QVBoxLayout(self.chat_widget_container)
        self.chat_layout.setAlignment(Qt.AlignmentFlag.AlignTop)
        self.chat_area.setWidget(self.chat_widget_container)
        right_layout.addWidget(self.chat_area)

        self.attachment_scroll = QScrollArea()
        self.attachment_scroll.setWidgetResizable(True)
        self.attachment_scroll.setMaximumHeight(120)
        self.attachment_scroll.setVisible(False)
        self.attachment_scroll.setStyleSheet("border: none; background-color: #252525;")

        self.attachment_container = QWidget()
        self.attachment_container.setStyleSheet("background-color: #252525;")
        self.attachment_grid = QHBoxLayout(self.attachment_container)
        self.attachment_grid.setAlignment(Qt.AlignmentFlag.AlignLeft)
        self.attachment_scroll.setWidget(self.attachment_container)
        right_layout.addWidget(self.attachment_scroll)

        input_layout = QHBoxLayout()
        self.attach_btn = QPushButton("📎")
        self.attach_btn.setFixedSize(40, 40)
        self.attach_btn.setStyleSheet("""
            QPushButton { background-color: #555; color: white; border-radius: 5px; font-size: 16px; }
            QPushButton:hover { background-color: #666; }
        """)
        self.attach_btn.clicked.connect(self.attach_file)

        self.input_field = QLineEdit()
        self.input_field.setPlaceholderText("Type a message...")
        self.input_field.setStyleSheet(
            "background: #3A3A3A; color: white; padding: 10px; border: 1px solid #555; border-radius: 5px;")
        self.input_field.returnPressed.connect(self.send_message)

        self.input_field.installEventFilter(self)

        self.send_button = QPushButton("Send")
        self.send_button.setStyleSheet("""
            QPushButton { background-color: #32CD32; color: white; padding: 10px 20px; font-weight: bold; border-radius: 5px; }
            QPushButton:hover { background-color: #2ECC71; }
        """)
        self.send_button.clicked.connect(self.handle_send_or_stop)

        input_layout.addWidget(self.attach_btn)
        input_layout.addWidget(self.input_field)
        input_layout.addWidget(self.send_button)
        right_layout.addLayout(input_layout)

        self.splitter.addWidget(self.sidebar_widget)
        self.splitter.addWidget(self.right_widget)
        self.splitter.setStretchFactor(1, 4)
        main_h_layout.addWidget(self.splitter)

        self.current_attachments = []
        self.current_streaming_bubble = None
        self.thread = None
        self.is_generating = False

        self.all_chats = ChatStorage.load_chats()
        self.current_chat_id = None
        self.filtered_chats = []

        self.refresh_sidebar()
        if not self.all_chats:
            self.create_new_chat()
        else:
            first_chat_id = list(self.all_chats.keys())[0]
            self.load_chat_by_id(first_chat_id)

        MDL2_ICONS = {
            'attach': '\uE723',
            'folder': '\uE8B7',
            'settings': '\uE713',
            'new': '\uE710',
            'search': '\uE721',
            'file': '\uE8A5',
            'image': '\uE91B',
            'video': '\uE714',
            'delete': '\uE74D',
        }

        self.attach_btn.setFont(QFont("Segoe MDL2 Assets", 16))
        self.attach_btn.setText(MDL2_ICONS['attach'])

        self.settings_btn.setFont(QFont("Segoe MDL2 Assets", 13))
        self.settings_btn.setText(f"{MDL2_ICONS['settings']} Settings")

        self.prompt_btn.setFont(QFont("Segoe MDL2 Assets", 13))
        self.prompt_btn.setText(f"{MDL2_ICONS['settings']} Settings")

        #self.workspace_btn.setFont(QFont("Segoe MDL2 Assets", 13))
        #self.workspace_btn.setText(f"{MDL2_ICONS['folder']} {WORKSPACE_FOLDER}")



    def open_settings(self):
        dialog = SettingsDialog(self)
        dialog.exec()

    def on_model_type_changed(self, model_type):
        self.model_combo.clear()
        if model_type == "Text Models":
            self.model_combo.addItems(TEXT_MODELS)
            self.input_field.setPlaceholderText("Type a message...")
        else:
            self.model_combo.addItems(IMAGE_MODELS)
            self.input_field.setPlaceholderText("Describe the image you want to generate...")

    def is_image_model(self):
        return self.model_type_combo.currentText() == "Image Models"

    def header_btn_style(self):
        return """
            QPushButton {
                background-color: #404040;
                color: #ffffff;
                border: 1px solid #555;
                border-radius: 6px;
                padding: 6px 15px;
                font-size: 13px;
            }
            QPushButton:hover {
                background-color: #4a4a4a;
                border: 1px solid #32CD32;
            }
        """

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

    def dropEvent(self, event):
        for url in event.mimeData().urls():
            path = url.toLocalFile()
            if os.path.exists(path):
                self.process_file_path(path)

    def eventFilter(self, source, event):
        if source == self.input_field and event.type() == QEvent.Type.KeyPress:
            if (event.modifiers() == Qt.KeyboardModifier.ControlModifier and
                    event.key() == Qt.Key.Key_V):

                clipboard = QApplication.clipboard()
                mime_data = clipboard.mimeData()

                if mime_data.hasUrls():
                    urls = mime_data.urls()
                    if urls:
                        path = urls[0].toLocalFile()
                        if os.path.exists(path):
                            self.process_file_path(path)
                            return True

                elif mime_data.hasImage():
                    image = clipboard.image()
                    if not image.isNull():
                        self.process_clipboard_image(image)
                        return True

        return super().eventFilter(source, event)

    def process_file_path(self, path):
        ext = os.path.splitext(path)[1].lower()
        attachment = {"path": os.path.basename(path), "content": None, "is_image": False, "is_video": False}

        if FileConverter.is_image(path):
            saved_path = FileManager.save_media_from_bytes(open(path, "rb").read(), ext)
            if saved_path:
                attachment["content"] = saved_path
                attachment["is_image"] = True
            else:
                QMessageBox.warning(self, "Error", "Could not process image.")
                return
        elif FileConverter.is_video(path):
            saved_path = FileManager.save_media_from_bytes(open(path, "rb").read(), ext)
            if saved_path:
                attachment["content"] = saved_path
                attachment["is_video"] = True
                thumb = FileConverter.get_video_thumbnail(os.path.join(WORKSPACE_FOLDER, saved_path))
                if thumb:
                    attachment["thumbnail"] = thumb
            else:
                QMessageBox.warning(self, "Error", "Could not process video.")
                return
        else:
            content = FileConverter.convert(path)
            if content.startswith("[Error"):
                QMessageBox.warning(self, "Conversion Error", content)
                return

            attachment["content"] = content
            attachment["is_image"] = False
            attachment["needs_chunking"] = len(content) > AUTO_CHUNK_THRESHOLD

        self.current_attachments.append(attachment)
        self.update_attachment_preview()

    def send_file_in_chunks_silent(self, content, filename):
        chunk_size = 30000

        if len(content) > 1000000:
            print(f"Warning: File {filename} too large ({len(content)} chars), truncating to 1MB")
            content = content[:1000000]

        chunks = [content[i:i + chunk_size] for i in range(0, len(content), chunk_size)]

        for idx, chunk in enumerate(chunks, 1):
            chunk_content = f"[ATTACHED FILE: {filename} - Part {idx}/{len(chunks)}]\n{chunk}\n[END FILE]"

            message_payload = {"role": "user", "content": chunk_content}
            self.all_chats[self.current_chat_id]["messages"].append(message_payload)

    def split_file_to_chunks(self, filepath, chunk_size=100000):
        try:
            content = FileConverter.convert(filepath)
            if content.startswith("[Error"):
                return None

            chunks = [content[i:i + chunk_size] for i in range(0, len(content), chunk_size)]
            base_name = os.path.splitext(os.path.basename(filepath))[0]

            chunk_files = []
            for idx, chunk in enumerate(chunks, 1):
                chunk_filename = f"{base_name}_chunk{idx}.txt"
                FileManager.write_file(chunk_filename, chunk)
                chunk_files.append(chunk_filename)

            return chunk_files
        except Exception as e:
            return None

    def send_file_in_chunks(self, filepath, content, filename):
        chunk_size = 30000
        chunks = [content[i:i + chunk_size] for i in range(0, len(content), chunk_size)]

        if not self.current_chat_id:
            self.create_new_chat()

        self.update_chat_title_if_needed(f"File: {filename}")

        for idx, chunk in enumerate(chunks, 1):
            chunk_content = f"[ATTACHED FILE: {filename}]\n{chunk}\n[END FILE]"

            message_payload = {"role": "user", "content": chunk_content}
            self.all_chats[self.current_chat_id]["messages"].append(message_payload)

        display_text = f"📄 [{filename}]"
        self.add_message_block("user", display_text)

        self.save_current_chat()

        QTimer.singleShot(100, lambda: self.chat_area.verticalScrollBar().setValue(
            self.chat_area.verticalScrollBar().maximum()))

    def process_clipboard_image(self, qimage):
        try:
            byte_array = QByteArray()
            buffer = QBuffer(byte_array)
            buffer.open(QIODevice.OpenModeFlag.WriteOnly)
            qimage.save(buffer, "PNG")
            image_bytes = byte_array.data()

            saved_path = FileManager.save_media_from_bytes(image_bytes, ".png")
            if saved_path:
                attachment = {
                    "path": "pasted_image.png",
                    "content": saved_path,
                    "is_image": True,
                    "is_video": False,
                    "qimage": qimage
                }
                self.current_attachments.append(attachment)
                self.update_attachment_preview()
            else:
                QMessageBox.warning(self, "Paste Error", "Failed to save pasted image.")
        except Exception as e:
            QMessageBox.warning(self, "Paste Error", f"Failed to paste image: {e}")

    def update_attachment_preview(self):
        while self.attachment_grid.count():
            child = self.attachment_grid.takeAt(0)
            if child.widget():
                child.widget().deleteLater()

        if not self.current_attachments:
            self.attachment_scroll.setVisible(False)
            return

        self.attachment_scroll.setVisible(True)

        for idx, attachment in enumerate(self.current_attachments):
            preview_widget = QWidget()
            preview_widget.setStyleSheet("background-color: #333; border-radius: 5px; padding: 5px;")
            preview_layout = QVBoxLayout(preview_widget)
            preview_layout.setContentsMargins(5, 5, 5, 5)
            preview_layout.setSpacing(3)

            if attachment.get("is_image", False):
                img_label = QLabel()
                img_label.setFixedSize(50, 50)
                img_label.setScaledContents(True)
                img_label.setStyleSheet("border: 1px solid #555; border-radius: 3px;")

                if "qimage" in attachment:
                    pixmap = QPixmap.fromImage(attachment["qimage"])
                else:
                    pixmap = QPixmap(os.path.join(WORKSPACE_FOLDER, attachment["content"]))

                img_label.setPixmap(pixmap.scaled(50, 50, Qt.AspectRatioMode.KeepAspectRatio,
                                                  Qt.TransformationMode.SmoothTransformation))
                preview_layout.addWidget(img_label)
            elif attachment.get("is_video", False):
                icon_label = QLabel()
                icon_label.setFixedSize(50, 50)
                icon_label.setScaledContents(True)
                icon_label.setStyleSheet("border: 1px solid #555; border-radius: 3px;")

                if "thumbnail" in attachment and attachment["thumbnail"]:
                    pixmap = attachment["thumbnail"]
                    overlay = QPixmap(pixmap.size())
                    overlay.fill(Qt.GlobalColor.transparent)
                    painter = QPainter(overlay)
                    painter.drawPixmap(0, 0, pixmap)
                    painter.setBrush(QColor(0, 0, 0, 128))
                    painter.setPen(Qt.PenStyle.NoPen)
                    painter.drawRect(overlay.rect())
                    painter.setPen(QColor(255, 255, 255))
                    font = painter.font()
                    font.setPointSize(24)
                    painter.setFont(font)
                    painter.drawText(overlay.rect(), Qt.AlignmentFlag.AlignCenter, "▶")
                    painter.end()
                    icon_label.setPixmap(overlay.scaled(50, 50, Qt.AspectRatioMode.KeepAspectRatio,
                                                        Qt.TransformationMode.SmoothTransformation))
                else:
                    icon_label.setText("🎥")
                    icon_label.setStyleSheet("font-size: 30px; border: 1px solid #555; border-radius: 3px;")
                    icon_label.setAlignment(Qt.AlignmentFlag.AlignCenter)

                preview_layout.addWidget(icon_label)
            else:
                file_icon = QLabel("📄")
                file_icon.setStyleSheet("font-size: 30px;")
                file_icon.setAlignment(Qt.AlignmentFlag.AlignCenter)
                file_icon.setFixedSize(50, 50)
                preview_layout.addWidget(file_icon)

            name_label = QLabel(attachment["path"])
            name_label.setStyleSheet("color: #aaa; font-size: 10px;")
            name_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            name_label.setWordWrap(True)
            name_label.setMaximumWidth(60)
            preview_layout.addWidget(name_label)

            remove_btn = QPushButton("×")
            remove_btn.setToolTip("Remove attachment")
            remove_btn.setCursor(Qt.CursorShape.PointingHandCursor)
            remove_btn.setFixedSize(18, 18)
            remove_btn.setStyleSheet("""
                QPushButton {
                    background-color: rgba(0, 0, 0, 0);
                    color: rgba(255, 255, 255, 0.65);
                    border: 1px solid rgba(255, 255, 255, 0.18);
                    border-radius: 9px;
                    font-weight: 700;
                    font-size: 13px;
                    padding: 0px;
                }
                QPushButton:hover {
                    background-color: rgba(231, 76, 60, 0.18);
                    border: 1px solid rgba(231, 76, 60, 0.55);
                    color: rgba(255, 255, 255, 0.95);
                }
                QPushButton:pressed {
                    background-color: rgba(231, 76, 60, 0.28);
                    border: 1px solid rgba(231, 76, 60, 0.75);
                }
            """)
            remove_btn.clicked.connect(lambda checked, i=idx: self.remove_attachment(i))
            preview_layout.addWidget(remove_btn, alignment=Qt.AlignmentFlag.AlignCenter)

            self.attachment_grid.addWidget(preview_widget)

        self.attachment_grid.addStretch()

    def remove_attachment(self, index):
        if 0 <= index < len(self.current_attachments):
            self.current_attachments.pop(index)
            self.update_attachment_preview()

    def filter_chats(self, text):
        text = text.lower()
        self.chat_list_widget.clear()

        sorted_chats = sorted(
            self.all_chats.items(),
            key=lambda item: item[1].get('timestamp', 0),
            reverse=True
        )

        for chat_id, chat_data in sorted_chats:
            name = chat_data.get('name', 'New Chat').lower()
            content_match = False
            if chat_data.get("messages"):
                last_msg = str(chat_data["messages"][-1].get("content", "")).lower()
                if text in last_msg:
                    content_match = True

            if text in name or content_match or text == "":
                item = QListWidgetItem(chat_data.get('name', 'New Chat'))
                item.setData(Qt.ItemDataRole.UserRole, chat_id)
                self.chat_list_widget.addItem(item)
                if chat_id == self.current_chat_id:
                    item.setSelected(True)

    def refresh_sidebar(self):
        self.filter_chats(self.search_input.text())

    def create_new_chat(self):
        new_id = str(uuid.uuid4())
        timestamp = time.time()
        self.all_chats[new_id] = {
            "name": "New Chat",
            "timestamp": timestamp,
            "messages": [{"role": "system", "content": DEFAULT_SYSTEM_PROMPT}]
        }
        ChatStorage.save_chats(self.all_chats)
        self.refresh_sidebar()
        self.load_chat_by_id(new_id)

    def load_selected_chat(self, item):
        chat_id = item.data(Qt.ItemDataRole.UserRole)
        if chat_id != self.current_chat_id:
            self.load_chat_by_id(chat_id)

    def load_chat_by_id(self, chat_id):
        self.current_chat_id = chat_id

        if self.is_generating:
            self.stop_generation()

        while self.chat_layout.count():
            child = self.chat_layout.takeAt(0)
            if child.widget():
                child.widget().deleteLater()

        messages = self.all_chats[chat_id]["messages"]

        skip_counter = 0

        for idx, msg in enumerate(messages):
            if skip_counter > 0:
                skip_counter -= 1
                continue

            try:
                content = msg["content"]
                role = msg["role"]

                if role == "system" and ":::create" not in str(content) and ":::delete" not in str(
                        content) and idx == 0:
                    continue

                if role == "user":
                    display_text = ""
                    media_data_list = []

                    str_content_for_check = ""
                    if isinstance(content, str):
                        str_content_for_check = content
                    elif isinstance(content, list):
                        for part in content:
                            if part.get("type") == "text":
                                str_content_for_check += part.get("text", "")
                            elif part.get("type") == "image_url":
                                url = part.get("image_url", {}).get("url", "")
                                if "base64," in url:
                                    media_data_list.append({"type": "image", "data": url.split("base64,")[1]})
                                elif not url.startswith("http"):
                                    full_path = os.path.join(WORKSPACE_FOLDER, url)
                                    if os.path.exists(full_path):
                                        b64 = FileConverter.encode_media_base64(full_path)
                                        if b64: media_data_list.append({"type": "image", "data": b64})
                            elif part.get("type") == "video_url":
                                url = part.get("video_url", {}).get("url", "")

                                if not url.startswith("data:") and not url.startswith("http"):
                                    full_path = os.path.join(WORKSPACE_FOLDER, url)
                                    if os.path.exists(full_path):
                                        media_data_list.append({"type": "video", "path": full_path})
                                        fname = os.path.basename(full_path)
                                        if f"🎥 [Video: {fname}]" not in str_content_for_check:
                                            pass
                                            #str_content_for_check += f"\n🎥 [Video: {fname}]"

                    if idx + 1 < len(messages):
                        next_msg = messages[idx + 1]
                        next_content = next_msg.get("content", "")
                        next_str = ""
                        if isinstance(next_content, str):
                            next_str = next_content
                        elif isinstance(next_content, list):
                            for p in next_content:
                                if p.get("type") == "text": next_str += p.get("text", "")

                        chunk_match = re.search(r"\[ATTACHED FILE: (.*?) - Part 1/(\d+)\]", next_str)
                        if chunk_match:
                            filename = chunk_match.group(1)
                            total_parts = int(chunk_match.group(2))
                            display_text = self._clean_file_tags(str_content_for_check)
                            display_text += f"\n\n📄 [{filename}]"
                            self.add_message_block("user", display_text, media_data_list=media_data_list)
                            skip_counter = total_parts
                            continue

                    chunk_match = re.search(r"\[ATTACHED FILE: (.*?) - Part (\d+)/(\d+)\]", str_content_for_check)
                    if chunk_match:
                        filename = chunk_match.group(1)
                        curr_p = int(chunk_match.group(2))
                        total_p = int(chunk_match.group(3))
                        if curr_p == 1:
                            self.add_message_block("user", f"📄 [{filename}]", media_data_list=media_data_list)
                            skip_counter = total_p - curr_p
                            continue

                    display_text = self._clean_file_tags(str_content_for_check)
                    self.add_message_block("user", display_text, media_data_list=media_data_list)

                elif role == "assistant":
                    if isinstance(content, list):
                        text_content = ""
                        has_images = False
                        for part in content:
                            if part.get("type") == "text":
                                text_content += part.get("text", "")
                            elif part.get("type") == "image_url":
                                has_images = True
                                url = part.get("image_url", {}).get("url", "")
                                b64 = self._get_b64_from_url(url)
                                if b64: self._add_assistant_image(b64)

                        if text_content.strip():
                            self.add_message_block("assistant", text_content)
                        if has_images: continue

                    elif "🖼️ Generated image:" in str(content):
                        image_path_match = re.search(r'images/[a-f0-9\-]+\.png', content)
                        if image_path_match:
                            full_path = os.path.join(WORKSPACE_FOLDER, image_path_match.group())
                            b64 = FileConverter.encode_media_base64(full_path)
                            if b64:
                                self._add_assistant_image(b64)
                                continue

                    if isinstance(content, str):
                        self.add_message_block("assistant", content)

                elif role == "system":
                    if "⚙️" in str(content) or "🗑️" in str(content):
                        self.add_message_block("system", content)

            except Exception as e:
                print(f"Error loading message: {e}")
                continue

        for i in range(self.chat_list_widget.count()):
            item = self.chat_list_widget.item(i)
            if item.data(Qt.ItemDataRole.UserRole) == chat_id:
                item.setSelected(True)
                break
        QTimer.singleShot(100, lambda: self.chat_area.verticalScrollBar().setValue(
            self.chat_area.verticalScrollBar().maximum()))

    def _clean_file_tags(self, text):
        if "[ATTACHED FILE:" not in text: return text
        lines = text.split("\n")
        clean = []
        in_file = False
        for line in lines:
            if line.startswith("[ATTACHED FILE:"):
                in_file = True
                fname = line.replace("[ATTACHED FILE:", "").replace("]", "").strip().split(" - Part")[0]
                clean.append(f"📄 [{fname}]")
            elif line.startswith("[END FILE]"):
                in_file = False
            elif not in_file:
                clean.append(line)
        return "\n".join(clean)

    def _get_b64_from_url(self, url):
        if "base64," in url: return url.split("base64,")[1]
        if not url.startswith("http"):
            full_path = os.path.join(WORKSPACE_FOLDER, url)
            if os.path.exists(full_path): return FileConverter.encode_media_base64(full_path)
        return None

    def _add_assistant_image(self, b64):
        container = QWidget()
        h_layout = QHBoxLayout(container)
        h_layout.setContentsMargins(0, 0, 0, 0)
        image_widget = self.create_image_preview({"type": "image", "data": b64}, "assistant")
        h_layout.addWidget(image_widget)
        h_layout.addStretch()
        self.chat_layout.addWidget(container)
        self.chat_layout.addSpacing(6)

    def reconstruct_chunked_file(self, filename, messages):
        """Helper to stitch together split file chunks from history"""
        full_content = []
        found_parts = {}
        total_parts = 0

        escaped_name = re.escape(filename)
        header_pattern = re.compile(r'\[ATTACHED FILE: ' + escaped_name + r' - Part (\d+)/(\d+)\]\n(.*?)\n\[END FILE\]',
                                    re.DOTALL)

        for msg in messages:
            content = msg.get("content", "")
            if isinstance(content, list):
                text_content = ""
                for part in content:
                    if part.get("type") == "text":
                        text_content += part.get("text", "")
                content = text_content

            match = header_pattern.search(str(content))
            if match:
                part_idx = int(match.group(1))
                total_parts = int(match.group(2))
                file_content = match.group(3)
                found_parts[part_idx] = file_content

        if found_parts:
            for i in range(1, total_parts + 1):
                if i in found_parts:
                    full_content.append(found_parts[i])
            return "".join(full_content)

        return None

    def find_file_content_in_history(self, filename):
        if not self.current_chat_id or self.current_chat_id not in self.all_chats:
            return None

        messages = self.all_chats[self.current_chat_id]["messages"]

        chunked_content = self.reconstruct_chunked_file(filename, messages)
        if chunked_content:
            return chunked_content

        for msg in reversed(messages):
            content = msg.get("content", "")

            if isinstance(content, list):
                text_content = ""
                for part in content:
                    if part.get("type") == "text":
                        text_content += part.get("text", "")
                content = text_content

            pattern = r'\[ATTACHED FILE: ' + re.escape(filename) + r'\]\n(.*?)\n\[END FILE\]'
            match = re.search(pattern, content, re.DOTALL)
            if match:
                return match.group(1)
        return None

    def save_current_chat(self):
        if self.current_chat_id and self.current_chat_id in self.all_chats:
            self.all_chats[self.current_chat_id]["timestamp"] = time.time()
            ChatStorage.save_chats(self.all_chats)

    def update_chat_title_if_needed(self, user_text):
        if not self.current_chat_id: return

        chat_data = self.all_chats[self.current_chat_id]
        if chat_data["name"] == "New Chat":
            new_name = user_text[:30] + "..." if len(user_text) > 30 else user_text
            chat_data["name"] = new_name
            self.refresh_sidebar()

    def show_context_menu(self, pos):
        item = self.chat_list_widget.itemAt(pos)
        if not item:
            return

        menu = QMenu(self)
        menu.setStyleSheet("""
            QMenu { background-color: #333; color: white; border: 1px solid #555; }
            QMenu::item { padding: 6px 20px; }
            QMenu::item:selected { background-color: #32CD32; }
        """)

        rename_action = menu.addAction("✏️ Rename")
        export_json_action = menu.addAction("💾 Export JSON")
        export_md_action = menu.addAction("📝 Export Markdown")
        delete_action = menu.addAction("🗑️ Delete")

        action = menu.exec(self.chat_list_widget.mapToGlobal(pos))

        if action == rename_action:
            self.rename_chat_item(item)
        elif action == delete_action:
            self.delete_chat_item(item)
        elif action == export_json_action:
            self.export_chat(item, "json")
        elif action == export_md_action:
            self.export_chat(item, "md")

    def export_chat(self, item, format_type):
        chat_id = item.data(Qt.ItemDataRole.UserRole)
        chat_data = self.all_chats.get(chat_id)
        if not chat_data: return

        filename = f"chat_export_{int(time.time())}"
        if format_type == "json":
            filename += ".json"
            save_path, _ = QFileDialog.getSaveFileName(self, "Export JSON", filename, "JSON Files (*.json)")
            if save_path:
                with open(save_path, 'w', encoding='utf-8') as f:
                    json.dump(chat_data, f, indent=4)
        else:
            filename += ".md"
            save_path, _ = QFileDialog.getSaveFileName(self, "Export Markdown", filename, "Markdown Files (*.md)")
            if save_path:
                with open(save_path, 'w', encoding='utf-8') as f:
                    f.write(f"# {chat_data.get('name')}\n\n")
                    for msg in chat_data.get('messages', []):
                        role = msg['role'].upper()
                        content = msg['content']
                        if isinstance(content, list): content = "[Complex Content]"
                        f.write(f"**{role}**: {content}\n\n---\n\n")

    def rename_chat_item(self, item):
        chat_id = item.data(Qt.ItemDataRole.UserRole)
        if not chat_id: return

        current_name = self.all_chats[chat_id]["name"]

        dialog = QInputDialog(self)
        dialog.setStyleSheet(DIALOG_STYLESHEET)
        dialog.setWindowTitle("Rename Chat")
        dialog.setLabelText("New Name:")
        dialog.setTextValue(current_name)
        dialog.setOkButtonText("Save")

        ok = dialog.exec()
        new_name = dialog.textValue()

        if ok and new_name.strip():
            self.all_chats[chat_id]["name"] = new_name.strip()
            ChatStorage.save_chats(self.all_chats)
            item.setText(new_name.strip())

    def delete_chat_item(self, item):
        chat_id = item.data(Qt.ItemDataRole.UserRole)
        if not chat_id: return

        msg_box = QMessageBox(self)
        msg_box.setStyleSheet(DIALOG_STYLESHEET)
        msg_box.setWindowTitle("Delete Chat")
        msg_box.setText("Are you sure you want to delete this chat permanently?")
        msg_box.setStandardButtons(QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
        msg_box.setDefaultButton(QMessageBox.StandardButton.No)

        confirm = msg_box.exec()

        if confirm == QMessageBox.StandardButton.Yes:
            if chat_id in self.all_chats:
                del self.all_chats[chat_id]
                ChatStorage.save_chats(self.all_chats)

            if self.current_chat_id == chat_id:
                self.current_chat_id = None
                while self.chat_layout.count():
                    child = self.chat_layout.takeAt(0)
                    if child.widget(): child.widget().deleteLater()

                if self.all_chats:
                    next_id = list(self.all_chats.keys())[0]
                    self.load_chat_by_id(next_id)
                else:
                    self.create_new_chat()

            self.refresh_sidebar()

    def open_workspace_folder(self):
        path = os.path.abspath(WORKSPACE_FOLDER)
        QDesktopServices.openUrl(QUrl.fromLocalFile(path))

    def open_prompt_editor(self):
        if not self.current_chat_id: return
        messages = self.all_chats[self.current_chat_id]["messages"]
        system_msg = next((m for m in messages if m["role"] == "system"), None)
        current_prompt = system_msg["content"] if system_msg else DEFAULT_SYSTEM_PROMPT

        dialog = PromptEditorDialog(current_prompt, self)
        if dialog.exec():
            new_prompt = dialog.get_prompt()
            if system_msg:
                system_msg["content"] = new_prompt
            else:
                messages.insert(0, {"role": "system", "content": new_prompt})
            self.save_current_chat()
            QMessageBox.information(self, "Info", "System prompt updated for this chat.")

    def attach_file(self):
        file_filter = "All Supported (*.txt *.md *.py *.js *.html *.css *.json *.pdf *.docx *.png *.jpg *.jpeg *.mp4 *.avi *.mov *.mkv *.webm);;All Files (*)"
        path, _ = QFileDialog.getOpenFileName(self, "Attach File", filter=file_filter)
        if path:
            self.process_file_path(path)

    def clear_attachment(self):
        self.current_attachments = []
        self.update_attachment_preview()

    def add_message_block(self, sender, text=None, streaming=False, media_data_list=None):
        if streaming and sender == "assistant":
            bubble = StreamingMessageBubble(sender)
            bubble.start_streaming()
            self.current_streaming_bubble = bubble
        else:
            bubble = StreamingMessageBubble(sender, text or "")

        bubble.label.clicked.connect(lambda: self.check_file_click(bubble.full_text))

        container = QWidget()
        v_layout = QVBoxLayout(container)
        v_layout.setContentsMargins(0, 0, 0, 0)
        v_layout.setSpacing(5)

        if media_data_list:
            for media_data in media_data_list:
                image_widget = self.create_image_preview(media_data, sender)
                if sender == "user":
                    image_container = QWidget()
                    h_img = QHBoxLayout(image_container)
                    h_img.setContentsMargins(0, 0, 0, 0)
                    h_img.addStretch()
                    h_img.addWidget(image_widget)
                    v_layout.addWidget(image_container)
                else:
                    v_layout.addWidget(image_widget)

        bubble_container = QWidget()
        h = QHBoxLayout(bubble_container)
        h.setContentsMargins(0, 0, 0, 0)
        if sender == "user":
            h.addStretch()
            h.addWidget(bubble)
        else:
            h.addWidget(bubble)
            h.addStretch()

        v_layout.addWidget(bubble_container)

        self.chat_layout.addWidget(container)
        self.chat_layout.addSpacing(6)

        QTimer.singleShot(10, lambda: self.chat_area.verticalScrollBar().setValue(
            self.chat_area.verticalScrollBar().maximum()))

        return bubble

    def find_file_content_in_history(self, filename):
        if not self.current_chat_id or self.current_chat_id not in self.all_chats:
            return None

        messages = self.all_chats[self.current_chat_id]["messages"]
        for msg in reversed(messages):
            content = msg.get("content", "")

            if isinstance(content, list):
                text_content = ""
                for part in content:
                    if part.get("type") == "text":
                        text_content += part.get("text", "")
                content = text_content

            pattern = r'\[ATTACHED FILE: ' + re.escape(filename) + r'\]\n(.*?)\n\[END FILE\]'
            match = re.search(pattern, content, re.DOTALL)
            if match:
                return match.group(1)
        return None

    def check_file_click(self, text):
        matches = re.findall(r'📄 \[(.*?)\]', text)
        if not matches: return

        if len(matches) == 1:
            fname = matches[0]
            content = self.find_file_content_in_history(fname)
            self.open_file_viewer(fname, content)
        else:
            menu = QMenu(self)
            menu.setStyleSheet(
                """QMenu{background-color:#333;color:white;}QMenu::item:selected{background-color:#32CD32;}""")
            for f in matches:
                action = menu.addAction(f)
                action.triggered.connect(
                    lambda checked, fn=f: self.open_file_viewer(fn, self.find_file_content_in_history(fn)))
            menu.exec(QCursor.pos())

    def open_file_viewer(self, filename, content=None):
        viewer = FileViewerDialog(filename, content, self)
        viewer.exec()

    def open_image_viewer(self, pixmap):
        viewer = ImageViewerDialog(pixmap, self)
        viewer.exec()

    def open_video_player(self, video_path):
        player_dialog = VideoPlayerDialog(video_path, self)
        player_dialog.exec()

    def handle_send_or_stop(self):
        if self.is_generating:
            self.stop_generation()
        else:
            self.send_message()

    def stop_generation(self):
        if self.thread and self.thread.isRunning():
            self.thread.stop()
            try:
                self.thread.response_chunk.disconnect()
                self.thread.response_complete.disconnect()
                self.thread.error_occurred.disconnect()
                self.thread.image_generated.disconnect()
            except:
                pass

        self.is_generating = False
        self.send_button.setText("Send")
        self.send_button.setStyleSheet("""
            QPushButton { background-color: #32CD32; color: white; padding: 10px 20px; font-weight: bold; border-radius: 5px; }
            QPushButton:hover { background-color: #2ECC71; }
        """)

        if self.current_streaming_bubble:
            self.current_streaming_bubble.stop_streaming()

            if not self.current_streaming_bubble.full_text.strip():
                try:
                    bubble_parent = self.current_streaming_bubble.parent()
                    if bubble_parent:
                        main_container = bubble_parent.parent()
                        if main_container:
                            main_container.deleteLater()
                            self.chat_layout.removeWidget(main_container)
                except Exception as e:
                    print(f"Error removing empty bubble: {e}")

            self.current_streaming_bubble = None

    def send_message(self):
        user_text = self.input_field.text().strip()
        if not user_text and not self.current_attachments:
            return

        if not self.current_chat_id:
            self.create_new_chat()

        self.update_chat_title_if_needed(user_text)

        display_text = user_text
        attached_media_data_for_display = []

        if self.current_attachments:
            chunked_files = []
            regular_attachments = []

            for attachment in self.current_attachments:
                if attachment.get("needs_chunking"):
                    chunked_files.append(attachment)
                else:
                    regular_attachments.append(attachment)

            content_parts = []
            file_text_parts = []

            if user_text:
                content_parts.append({"type": "text", "text": user_text})

            for attachment in regular_attachments:
                if attachment.get("is_image"):
                    if display_text:
                        display_text += f"\n"

                    full_path = os.path.join(WORKSPACE_FOLDER, attachment["content"])
                    b64 = FileConverter.encode_media_base64(full_path)
                    if b64:
                        attached_media_data_for_display.append({"type": "image", "data": b64})

                    content_parts.append({
                        "type": "image_url",
                        "image_url": {
                            "url": attachment["content"]
                        }
                    })
                elif attachment.get("is_video"):
                    if display_text:
                        display_text += f"\n"

                    full_path = os.path.join(WORKSPACE_FOLDER, attachment["content"])
                    attached_media_data_for_display.append({"type": "video", "path": full_path})

                    content_parts.append({
                        "type": "video_url",
                        "video_url": {
                            "url": attachment["content"]
                        }
                    })

                else:
                    if display_text:
                        display_text += f"\n📄 [{attachment['path']}]"
                    else:
                        display_text = f"📄 [{attachment['path']}]"

                    file_text_parts.append(
                        f"\n\n[ATTACHED FILE: {attachment['path']}]\n{attachment['content']}\n[END FILE]")

            for attachment in chunked_files:
                if display_text:
                    display_text += f"\n📄 [{attachment['path']}]"
                else:
                    display_text = f"📄 [{attachment['path']}]"

            if regular_attachments:
                if any(a.get("is_image") or a.get("is_video") for a in regular_attachments):
                    if file_text_parts:
                        if not user_text:
                            content_parts.insert(0, {"type": "text", "text": "".join(file_text_parts)})
                        else:
                            content_parts[0]["text"] += "".join(file_text_parts)

                    message_payload = {
                        "role": "user",
                        "content": content_parts
                    }
                else:
                    full_prompt = (user_text if user_text else "") + "".join(file_text_parts)
                    message_payload = {"role": "user", "content": full_prompt}

                msgs = self.all_chats[self.current_chat_id]["messages"]
                msgs.append(message_payload)
            elif user_text:
                message_payload = {"role": "user", "content": user_text}
                msgs = self.all_chats[self.current_chat_id]["messages"]
                msgs.append(message_payload)

            for attachment in chunked_files:
                self.send_file_in_chunks_silent(attachment["content"], attachment["path"])

            self.clear_attachment()
        else:
            message_payload = {"role": "user", "content": user_text}
            msgs = self.all_chats[self.current_chat_id]["messages"]
            msgs.append(message_payload)

        self.input_field.clear()

        self.add_message_block("user", display_text, media_data_list=attached_media_data_for_display)

        self.save_current_chat()
        self.trigger_assistant_response()

    def trigger_assistant_response(self):
        self.is_generating = True
        self.send_button.setText("Stop")
        self.send_button.setStyleSheet("""
            QPushButton { background-color: #e74c3c; color: white; padding: 10px 20px; font-weight: bold; border-radius: 5px; }
            QPushButton:hover { background-color: #c0392b; }
        """)

        model = self.model_combo.currentText()
        messages_history = self.all_chats[self.current_chat_id]["messages"]

        is_image = self.is_image_model()

        if not is_image:
            self.add_message_block("assistant", streaming=True)

        self.thread = AIResponseThread(model, messages_history, is_image)
        self.thread.response_chunk.connect(self.handle_stream_chunk)
        self.thread.response_complete.connect(self.handle_stream_complete)
        self.thread.error_occurred.connect(self.display_error)
        self.thread.image_generated.connect(self.handle_image_generated)
        self.thread.start()

    def handle_image_generated(self, image_bytes, saved_path):

        base64_str = base64.b64encode(image_bytes).decode('utf-8')

        image_message = {
            "role": "assistant",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {
                        "url": saved_path
                    }
                },
                {
                    "type": "text",
                    "text": "Here is the generated image:"
                }
            ]
        }
        self.all_chats[self.current_chat_id]["messages"].append(image_message)

        container = QWidget()
        h_layout = QHBoxLayout(container)
        h_layout.setContentsMargins(0, 0, 0, 0)
        h_layout.setSpacing(0)

        image_widget = self.create_image_preview({"type": "image", "data": base64_str}, "assistant")
        h_layout.addWidget(image_widget)
        h_layout.addStretch()

        self.chat_layout.addWidget(container)
        self.chat_layout.addSpacing(6)

        QTimer.singleShot(10, lambda: self.chat_area.verticalScrollBar().setValue(
            self.chat_area.verticalScrollBar().maximum()))

        self.save_current_chat()

    def create_image_preview(self, media_data, sender):
        try:
            pixmap = QPixmap()
            is_video = False

            if media_data.get("type") == "image":
                image_bytes = base64.b64decode(media_data["data"])
                pixmap.loadFromData(image_bytes)
            elif media_data.get("type") == "video":
                is_video = True
                path = media_data["path"]
                thumb = FileConverter.get_video_thumbnail(path)
                if thumb:
                    pixmap = thumb
                else:
                    pixmap = QPixmap(250, 150)
                    pixmap.fill(QColor("#333"))

            if pixmap.isNull():
                error_label = QLabel("⚠️ Media load failed")
                error_label.setStyleSheet("color: #ff6b6b; padding: 10px;")
                return error_label

            image_label = ClickableImageLabel()
            image_label.setCursor(Qt.CursorShape.PointingHandCursor)

            scaled_pixmap = pixmap.scaled(250, 250, Qt.AspectRatioMode.KeepAspectRatio,
                                          Qt.TransformationMode.SmoothTransformation)

            if is_video:
                overlay = QPixmap(scaled_pixmap.size())
                overlay.fill(Qt.GlobalColor.transparent)
                painter = QPainter(overlay)
                painter.drawPixmap(0, 0, scaled_pixmap)

                painter.setBrush(QColor(0, 0, 0, 100))
                painter.setPen(Qt.PenStyle.NoPen)
                painter.drawRect(overlay.rect())

                painter.setBrush(QColor(255, 255, 255, 200))
                center = overlay.rect().center()
                size = 30
                triangle = QPolygonF([
                    QPointF(center.x() - size / 2, center.y() - size / 2),
                    QPointF(center.x() - size / 2, center.y() + size / 2),
                    QPointF(center.x() + size / 2, center.y())
                ])
                painter.drawPolygon(triangle)
                painter.end()
                image_label.setPixmap(overlay)
            else:
                image_label.setPixmap(scaled_pixmap)

            image_label.setStyleSheet("border-radius: 8px; background-color: #333; border: 1px solid #444;")

            if is_video:
                image_label.clicked.connect(lambda: self.open_video_player(media_data["path"]))
            else:
                image_label.clicked.connect(lambda: self.open_image_viewer(pixmap))

            return image_label
        except Exception as e:
            print(f"Error displaying media: {e}")
            error_label = QLabel("⚠️ Media error")
            error_label.setStyleSheet("color: #ff6b6b; padding: 10px;")
            return error_label

    def handle_stream_chunk(self, chunk):
        if self.current_streaming_bubble:
            sb = self.chat_area.verticalScrollBar()
            was_at_bottom = sb.value() >= (sb.maximum() - 30)

            self.current_streaming_bubble.add_text_chunk(chunk)
            if was_at_bottom:
                sb.setValue(sb.maximum())

    def handle_stream_complete(self, full_response):
        self.is_generating = False
        self.send_button.setText("Send")
        self.send_button.setStyleSheet("""
            QPushButton { background-color: #32CD32; color: white; padding: 10px 20px; font-weight: bold; border-radius: 5px; }
            QPushButton:hover { background-color: #2ECC71; }
        """)

        if not self.is_image_model():
            self.all_chats[self.current_chat_id]["messages"].append({
                "role": "assistant",
                "content": full_response
            })

            if self.current_streaming_bubble:
                self.current_streaming_bubble.set_complete_text(full_response)
                self.current_streaming_bubble.stop_streaming()
                self.current_streaming_bubble = None

        self.save_current_chat()
        self.execute_ai_file_commands(full_response if not self.is_image_model() else "")

    def execute_ai_file_commands(self, text):
        create_pattern = re.compile(r':::create file="(.*?)"(.*?):::(.*?):::end_create:::', re.DOTALL)
        creates = create_pattern.findall(text)

        for filename, extra_args, content in creates:
            filename = filename.strip()
            if content.startswith('\n'): content = content[1:]
            result_msg = FileManager.write_file(filename, content)

            sys_msg = f"⚙️ {result_msg}"

            self.all_chats[self.current_chat_id]["messages"].append({"role": "system", "content": sys_msg})

            self.add_message_block("system", sys_msg)

        delete_pattern = re.compile(r':::delete file="(.*?)"(.*?):::')
        deletes = delete_pattern.findall(text)

        for filename, extra_args in deletes:
            filename = filename.strip()
            result_msg = FileManager.delete_file(filename)

            sys_msg = f"🗑️ {result_msg}"
            self.all_chats[self.current_chat_id]["messages"].append({"role": "system", "content": sys_msg})

            self.add_message_block("system", sys_msg)

        if creates or deletes:
            self.save_current_chat()

    def display_error(self, error):
        self.is_generating = False
        self.send_button.setText("Send")
        self.send_button.setStyleSheet("""
            QPushButton { background-color: #32CD32; color: white; padding: 10px 20px; font-weight: bold; border-radius: 5px; }
            QPushButton:hover { background-color: #2ECC71; }
        """)

        if self.current_streaming_bubble:
            self.current_streaming_bubble.stop_streaming()
            self.current_streaming_bubble.set_complete_text(f"Error: {error}")
            self.current_streaming_bubble.setStyleSheet("background-color: #c0392b; border-radius: 12px; padding: 6px;")
            self.current_streaming_bubble = None

    def closeEvent(self, event):
        if self.thread and self.thread.isRunning():
            self.thread.terminate()
        event.accept()


if __name__ == "__main__":
    app = QApplication(sys.argv)

    mono_font = QFont("Segoe UI Symbol", 10)
    app.setFont(mono_font)

    window = AIChatApp()
    window.show()
    sys.exit(app.exec())
